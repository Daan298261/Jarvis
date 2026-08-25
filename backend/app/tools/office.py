from __future__ import annotations

import csv
import io
import json
import os
import platform
from contextlib import contextmanager
from pathlib import Path
from typing import Any, Iterable

from .base import RiskLevel, Tool, ToolResult
from .safety import resolve_allowed_path


def office_com_available() -> bool:
    if platform.system() != "Windows":
        return False
    try:
        import pythoncom  # noqa: F401
        import win32com.client  # noqa: F401
    except Exception:
        return False
    return True


def office_library_available(app: str | None = None) -> bool:
    wanted = {
        "word": "docx",
        "excel": "openpyxl",
        "powerpoint": "pptx",
    }
    if app:
        module = wanted.get(app)
        if not module:
            return False
        return _module_ok(module)
    return any(_module_ok(name) for name in wanted.values())


def _module_ok(name: str) -> bool:
    import importlib.util

    return importlib.util.find_spec(name) is not None

_PROGID = {
    "word": "Word.Application",
    "excel": "Excel.Application",
    "powerpoint": "PowerPoint.Application",
}


def office_runtime_available() -> bool:
    if platform.system() != "Windows":
        return False
    try:
        import win32com.client  # noqa: F401
        import pythoncom  # noqa: F401
    except Exception:
        return False
    return True


def office_available() -> tuple[bool, str]:
    if platform.system() != "Windows":
        return False, "Office COM is only available on Windows"
    try:
        import win32com.client  # noqa: F401
    except Exception:
        return False, "pywin32 is not installed"
    roots = [
        Path(os.environ.get("PROGRAMFILES", r"C:\Program Files")) / "Microsoft Office",
        Path(os.environ.get("PROGRAMFILES(X86)", r"C:\Program Files (x86)")) / "Microsoft Office",
        Path(os.environ.get("PROGRAMFILES", r"C:\Program Files")) / "Microsoft Office 16",
    ]
    if any(root.exists() for root in roots):
        return True, "Microsoft Office appears to be installed"
    return False, "Microsoft Office does not appear to be installed"


def _dispatch(progid: str):
    import pythoncom
    import win32com.client

    pythoncom.CoInitialize()
    return win32com.client.Dispatch(progid)


@contextmanager
def _com_app(progid: str, visible: bool = False):
    app = _dispatch(progid)
    try:
        if hasattr(app, "Visible"):
            app.Visible = visible
        if hasattr(app, "DisplayAlerts"):
            app.DisplayAlerts = False
        yield app
    finally:
        try:
            app.Quit()
        except Exception:
            pass


def _parse_table(content: str) -> list[list[str]]:
    text = (content or "").strip()
    if not text:
        return []
    if text[0] in "[{":
        try:
            data = json.loads(text)
        except json.JSONDecodeError:
            data = None
        if isinstance(data, list) and data and isinstance(data[0], dict):
            keys: list[str] = []
            for row in data:
                for key in row:
                    if key not in keys:
                        keys.append(str(key))
            return [keys] + [["" if row.get(k) is None else str(row.get(k, "")) for k in keys] for row in data]
        if isinstance(data, list):
            rows: list[list[str]] = []
            for row in data:
                if isinstance(row, list):
                    rows.append(["" if c is None else str(c) for c in row])
                else:
                    rows.append([str(row)])
            return rows
    dialect = csv.excel_tab if "\t" in text.splitlines()[0] else csv.excel
    reader = csv.reader(io.StringIO(text), dialect=dialect)
    return [[cell for cell in row] for row in reader]


def _table_text(rows: Iterable[Iterable[Any]]) -> str:
    lines = []
    for row in rows:
        cells = [("" if c is None else str(c)) for c in row]
        lines.append("\t".join(cells))
    return "\n".join(lines)


def _slides_from_content(content: str) -> list[str]:
    text = content or ""
    parts = [part.strip() for part in text.split("\n---\n")]
    return [part for part in parts if part] or ([text.strip()] if text.strip() else [""])


class OfficeTool(Tool):
    name = "office"
    description = (
        "Read, create, and edit Microsoft Word, Excel, and PowerPoint files. "
        "Uses Windows COM when Office is installed; otherwise uses python-docx / openpyxl / python-pptx. "
        "Actions: create, read, write, save_as, append, info. Always write to a new file unless the user "
        "asked for in-place edits. Paths must be inside allowed directories."
    )
    risk = RiskLevel.MEDIUM
    parameters = {
        "type": "object",
        "properties": {
            "app": {"type": "string", "enum": ["word", "excel", "powerpoint"]},
            "action": {"type": "string", "enum": ["create", "read", "write", "save_as", "append", "info"]},
            "path": {"type": "string"},
            "content": {"type": "string", "description": "Plain text, TSV/CSV/JSON table, or slides separated by ---"},
            "destination": {"type": "string"},
            "sheet": {"type": "string", "description": "Excel sheet name"},
            "backend": {"type": "string", "enum": ["auto", "com", "library"], "default": "auto"},
        },
        "required": ["app", "action"],
    }

    def __init__(self, context_getter=None) -> None:
        self.context_getter = context_getter or (lambda: {})

    def _resolve(self, path: str | None) -> Path:
        if not path:
            raise ValueError("path is required")
        allowed = list((self.context_getter() or {}).get("allowed_directories") or [])
        return resolve_allowed_path(path, allowed)

    def _info(self, app: str, path: str | None) -> ToolResult:
        progid = _PROGID.get(app, "Word.Application")
        lib_ok = office_library_available(app) if app else office_library_available()
        bits = [
            f"app={app or 'unspecified'}",
            f"os={platform.system()}",
            f"available={lib_ok or office_com_available()}",
            f"library_available={lib_ok}",
            f"com_available={office_com_available()}",
            f"progid={progid}",
            f"{progid}",
        ]
        data: dict[str, Any] = {
            "windows": platform.system() == "Windows",
            "available": lib_ok or office_com_available(),
            "backend": "none",
        }
        if path:
            target = Path(path)
            if not target.exists():
                return ToolResult(False, "", error=f"Office file not found: {path}")
            stat = target.stat()
            bits.extend(
                [
                    f"path={target.resolve()}",
                    f"size_bytes={stat.st_size}",
                    f"suffix={target.suffix}",
                ]
            )
        bits.append("COM was not launched.")
        if platform.system() != "Windows":
            bits.append("Office COM is only available on Windows. COM was not launched.")
        return ToolResult(True, "\n".join(bits), data=data)

    async def execute(self, **kwargs: Any) -> ToolResult:
        app = (kwargs.get("app") or "").lower()
        action = (kwargs.get("action") or "").lower()
        backend = (kwargs.get("backend") or "auto").lower()
        if app not in {"word", "excel", "powerpoint"}:
            return ToolResult(False, "", error="app must be word, excel, or powerpoint")
        if action not in {"create", "read", "write", "save_as", "append", "info"}:
            return ToolResult(False, "", error=f"Unknown action {action}")
        try:
            if action == "info" and backend != "com":
                if kwargs.get("path") and office_library_available(app):
                    try:
                        return self._run_library(app, "info", kwargs)
                    except (PermissionError, FileNotFoundError, ValueError):
                        raise
                    except Exception:
                        pass
                return self._info(app, kwargs.get("path"))
            chosen = self._choose_backend(app, backend)
            if chosen == "library":
                return self._run_library(app, action, kwargs)
            if action == "info":
                return self._info(app, kwargs.get("path"))
            return self._run_com(app, action, kwargs)
        except PermissionError as exc:
            return ToolResult(False, "", error=str(exc))
        except Exception as exc:
            message = str(exc)
            lowered = message.lower()
            if "com is not available" in lowered or "only available on windows" in lowered:
                return ToolResult(False, "", error=message if "COM is not available" in message else f"Office COM is not available: {message}")
            if "unavailable" in lowered:
                return ToolResult(False, "", error=message)
            return ToolResult(False, "", error=f"Office automation unavailable: {exc}")

    def _choose_backend(self, app: str, backend: str) -> str:
        com_ok = office_com_available()
        lib_ok = office_library_available(app)
        if backend == "com":
            if not com_ok:
                raise RuntimeError("Office COM is not available on this machine")
            return "com"
        if backend == "library":
            if not lib_ok:
                raise RuntimeError(f"Office library backend is not installed for {app}")
            return "library"
        if lib_ok:
            return "library"
        if com_ok:
            return "com"
        raise RuntimeError(
            "Office is unavailable. Install Microsoft Office (Windows COM) or "
            "python-docx / openpyxl / python-pptx"
        )

    def _run_library(self, app: str, action: str, kwargs: dict[str, Any]) -> ToolResult:
        if app == "word":
            return self._word_lib(action, kwargs)
        if app == "excel":
            return self._excel_lib(action, kwargs)
        return self._ppt_lib(action, kwargs)

    def _target(self, kwargs: dict[str, Any], *, creating: bool) -> Path:
        dest = kwargs.get("destination") or kwargs.get("path")
        if creating:
            path = self._resolve(dest)
            path.parent.mkdir(parents=True, exist_ok=True)
            return path
        return self._resolve(kwargs.get("path") or dest)

    def _word_lib(self, action: str, kwargs: dict[str, Any]) -> ToolResult:
        from docx import Document

        if action == "create":
            dest = self._target(kwargs, creating=True)
            doc = Document()
            text = kwargs.get("content") or ""
            if text:
                for para in text.splitlines() or [""]:
                    doc.add_paragraph(para)
            else:
                doc.add_paragraph("")
            doc.save(str(dest))
            return ToolResult(True, f"Created Word document {dest}", data={"path": str(dest), "backend": "library"})
        src = self._target(kwargs, creating=False)
        if not src.exists():
            return ToolResult(False, "", error=f"File not found: {src}")
        doc = Document(str(src))
        text = "\n".join(p.text for p in doc.paragraphs)
        if action == "read":
            return ToolResult(True, text, data={"path": str(src), "backend": "library", "paragraphs": len(doc.paragraphs)})
        if action == "info":
            return ToolResult(
                True,
                f"path={src}\nbackend=library\nparagraphs={len(doc.paragraphs)}\ncharacters={len(text)}",
                data={"path": str(src), "paragraphs": len(doc.paragraphs), "backend": "library"},
            )
        dest = self._resolve(kwargs.get("destination") or str(src)) if action == "save_as" else src
        dest.parent.mkdir(parents=True, exist_ok=True)
        if action in {"write", "save_as"}:
            fresh = Document()
            body = kwargs.get("content")
            if body is None:
                for para in doc.paragraphs:
                    fresh.add_paragraph(para.text)
            else:
                for para in (body.splitlines() or [""]):
                    fresh.add_paragraph(para)
            fresh.save(str(dest))
            return ToolResult(True, f"Saved {dest}", data={"path": str(dest), "backend": "library"})
        if action == "append":
            extra = kwargs.get("content") or ""
            if extra:
                for para in extra.splitlines() or [extra]:
                    doc.add_paragraph(para)
            dest = self._resolve(kwargs.get("destination") or str(src))
            dest.parent.mkdir(parents=True, exist_ok=True)
            doc.save(str(dest))
            return ToolResult(True, f"Appended and saved {dest}", data={"path": str(dest), "backend": "library"})
        return ToolResult(False, "", error=f"Unknown action {action}")

    def _excel_lib(self, action: str, kwargs: dict[str, Any]) -> ToolResult:
        from openpyxl import Workbook, load_workbook

        sheet_name = kwargs.get("sheet")
        if action == "create":
            dest = self._target(kwargs, creating=True)
            wb = Workbook()
            ws = wb.active
            if sheet_name:
                ws.title = sheet_name
            for r_idx, row in enumerate(_parse_table(kwargs.get("content") or ""), start=1):
                for c_idx, value in enumerate(row, start=1):
                    ws.cell(r_idx, c_idx, value)
            wb.save(str(dest))
            return ToolResult(True, f"Created workbook {dest}", data={"path": str(dest), "backend": "library"})
        src = self._target(kwargs, creating=False)
        if not src.exists():
            return ToolResult(False, "", error=f"File not found: {src}")
        wb = load_workbook(str(src))
        ws = wb[sheet_name] if sheet_name and sheet_name in wb.sheetnames else wb.active
        rows = [[cell.value for cell in row] for row in ws.iter_rows()]
        if action == "read":
            return ToolResult(True, _table_text(rows), data={"path": str(src), "sheet": ws.title, "backend": "library"})
        if action == "info":
            return ToolResult(
                True,
                f"path={src}\nbackend=library\nsheets={', '.join(wb.sheetnames)}\nactive={ws.title}\nrows={len(rows)}",
                data={"path": str(src), "sheets": wb.sheetnames, "backend": "library"},
            )
        dest = self._resolve(kwargs.get("destination") or str(src)) if action in {"write", "save_as", "append"} else src
        dest.parent.mkdir(parents=True, exist_ok=True)
        if action in {"write", "save_as"}:
            table = _parse_table(kwargs.get("content") or "")
            if action == "save_as" and kwargs.get("content") is None:
                wb.save(str(dest))
                return ToolResult(True, f"Saved {dest}", data={"path": str(dest), "backend": "library"})
            out = Workbook()
            out_ws = out.active
            if sheet_name:
                out_ws.title = sheet_name
            for r_idx, row in enumerate(table, start=1):
                for c_idx, value in enumerate(row, start=1):
                    out_ws.cell(r_idx, c_idx, value)
            out.save(str(dest))
            return ToolResult(True, f"Saved {dest}", data={"path": str(dest), "backend": "library"})
        if action == "append":
            table = _parse_table(kwargs.get("content") or "")
            start = ws.max_row + 1 if ws.max_row else 1
            if ws.max_row == 1 and all(c.value is None for c in ws[1]):
                start = 1
            for r_idx, row in enumerate(table, start=start):
                for c_idx, value in enumerate(row, start=1):
                    ws.cell(r_idx, c_idx, value)
            wb.save(str(dest))
            return ToolResult(True, f"Appended and saved {dest}", data={"path": str(dest), "backend": "library"})
        return ToolResult(False, "", error=f"Unknown action {action}")

    def _ppt_lib(self, action: str, kwargs: dict[str, Any]) -> ToolResult:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        if action == "create":
            dest = self._target(kwargs, creating=True)
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            for block in _slides_from_content(kwargs.get("content") or "Slide 1"):
                self._add_slide(prs, block)
            prs.save(str(dest))
            return ToolResult(True, f"Created presentation {dest}", data={"path": str(dest), "backend": "library"})
        src = self._target(kwargs, creating=False)
        if not src.exists():
            return ToolResult(False, "", error=f"File not found: {src}")
        prs = Presentation(str(src))
        slides = [self._slide_text(slide) for slide in prs.slides]
        if action == "read":
            return ToolResult(True, "\n---\n".join(slides), data={"path": str(src), "slides": len(slides), "backend": "library"})
        if action == "info":
            return ToolResult(
                True,
                f"path={src}\nbackend=library\nslides={len(slides)}",
                data={"path": str(src), "slides": len(slides), "backend": "library"},
            )
        dest = self._resolve(kwargs.get("destination") or str(src))
        dest.parent.mkdir(parents=True, exist_ok=True)
        if action in {"write", "save_as"}:
            if kwargs.get("content") is None and action == "save_as":
                prs.save(str(dest))
                return ToolResult(True, f"Saved {dest}", data={"path": str(dest), "backend": "library"})
            out = Presentation()
            out.slide_width = Inches(13.333)
            out.slide_height = Inches(7.5)
            for block in _slides_from_content(kwargs.get("content") or ""):
                self._add_slide(out, block)
            out.save(str(dest))
            return ToolResult(True, f"Saved {dest}", data={"path": str(dest), "backend": "library"})
        if action == "append":
            for block in _slides_from_content(kwargs.get("content") or ""):
                self._add_slide(prs, block)
            prs.save(str(dest))
            return ToolResult(True, f"Appended and saved {dest}", data={"path": str(dest), "backend": "library"})
        return ToolResult(False, "", error=f"Unknown action {action}")

    def _add_slide(self, prs: Any, block: str) -> None:
        from pptx.util import Inches, Pt

        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12), Inches(6.5))
        tf = box.text_frame
        tf.word_wrap = True
        lines = block.splitlines() or [""]
        tf.text = lines[0]
        for line in lines[1:]:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(18)

    def _slide_text(self, slide: Any) -> str:
        chunks: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                chunks.append(shape.text)
        return "\n".join(chunks)

    def _run_com(self, app: str, action: str, kwargs: dict[str, Any]) -> ToolResult:
        if app == "word":
            return self._word_com(action, kwargs)
        if app == "excel":
            return self._excel_com(action, kwargs)
        return self._ppt_com(action, kwargs)

    def _word_com(self, action: str, kwargs: dict[str, Any]) -> ToolResult:
        with _com_app("Word.Application") as word:
            if action == "create":
                dest = self._target(kwargs, creating=True)
                doc = word.Documents.Add()
                if kwargs.get("content"):
                    doc.Range().Text = kwargs["content"]
                doc.SaveAs(str(dest))
                doc.Close()
                return ToolResult(True, f"Created Word document {dest}", data={"path": str(dest), "backend": "com"})
            src = self._target(kwargs, creating=False)
            doc = word.Documents.Open(str(src))
            try:
                text = doc.Range().Text
                if action == "read":
                    return ToolResult(True, text, data={"path": str(src), "backend": "com"})
                if action == "info":
                    return ToolResult(
                        True,
                        f"path={src}\nbackend=com\ncharacters={len(text or '')}",
                        data={"path": str(src), "backend": "com"},
                    )
                dest = self._resolve(kwargs.get("destination") or str(src))
                dest.parent.mkdir(parents=True, exist_ok=True)
                if action == "append" and kwargs.get("content"):
                    doc.Range(doc.Content.End - 1, doc.Content.End - 1).Text = kwargs["content"]
                elif action in {"write", "save_as"} and kwargs.get("content") is not None:
                    doc.Range().Text = kwargs["content"]
                doc.SaveAs(str(dest))
                return ToolResult(True, f"Saved {dest}", data={"path": str(dest), "backend": "com"})
            finally:
                doc.Close(False)

    def _excel_com(self, action: str, kwargs: dict[str, Any]) -> ToolResult:
        with _com_app("Excel.Application") as excel:
            if action == "create":
                dest = self._target(kwargs, creating=True)
                wb = excel.Workbooks.Add()
                sheet = wb.Worksheets(1)
                if kwargs.get("sheet"):
                    sheet.Name = kwargs["sheet"]
                self._write_com_sheet(sheet, _parse_table(kwargs.get("content") or ""))
                wb.SaveAs(str(dest))
                wb.Close()
                return ToolResult(True, f"Created workbook {dest}", data={"path": str(dest), "backend": "com"})
            src = self._target(kwargs, creating=False)
            wb = excel.Workbooks.Open(str(src))
            try:
                sheet = wb.Worksheets(kwargs["sheet"]) if kwargs.get("sheet") else wb.Worksheets(1)
                used = sheet.UsedRange
                rows = []
                values = used.Value
                if values is None:
                    rows = []
                elif not isinstance(values, tuple):
                    rows = [[values]]
                else:
                    for row in values:
                        if isinstance(row, tuple):
                            rows.append(["" if c is None else str(c) for c in row])
                        else:
                            rows.append(["" if row is None else str(row)])
                if action == "read":
                    return ToolResult(True, _table_text(rows), data={"path": str(src), "backend": "com"})
                if action == "info":
                    names = [wb.Worksheets(i).Name for i in range(1, wb.Worksheets.Count + 1)]
                    return ToolResult(
                        True,
                        f"path={src}\nbackend=com\nsheets={', '.join(names)}\nrows={len(rows)}",
                        data={"path": str(src), "sheets": names, "backend": "com"},
                    )
                dest = self._resolve(kwargs.get("destination") or str(src))
                dest.parent.mkdir(parents=True, exist_ok=True)
                table = _parse_table(kwargs.get("content") or "")
                if action == "append":
                    start = (used.Row + used.Rows.Count) if values is not None else 1
                    self._write_com_sheet(sheet, table, start_row=start)
                elif action in {"write", "save_as"} and kwargs.get("content") is not None:
                    sheet.Cells.Clear()
                    self._write_com_sheet(sheet, table)
                wb.SaveAs(str(dest))
                return ToolResult(True, f"Saved {dest}", data={"path": str(dest), "backend": "com"})
            finally:
                wb.Close(False)

    def _ppt_com(self, action: str, kwargs: dict[str, Any]) -> ToolResult:
        with _com_app("PowerPoint.Application") as ppt:
            if action == "create":
                dest = self._target(kwargs, creating=True)
                pres = ppt.Presentations.Add()
                for block in _slides_from_content(kwargs.get("content") or "Slide 1"):
                    slide = pres.Slides.Add(pres.Slides.Count + 1, 12)
                    try:
                        slide.Shapes.Title.TextFrame.TextRange.Text = block.splitlines()[0]
                    except Exception:
                        pass
                pres.SaveAs(str(dest))
                pres.Close()
                return ToolResult(True, f"Created presentation {dest}", data={"path": str(dest), "backend": "com"})
            src = self._target(kwargs, creating=False)
            pres = ppt.Presentations.Open(str(src), WithWindow=False)
            try:
                texts = []
                for i in range(1, pres.Slides.Count + 1):
                    slide = pres.Slides(i)
                    chunks = []
                    for shape in slide.Shapes:
                        try:
                            chunks.append(shape.TextFrame.TextRange.Text)
                        except Exception:
                            continue
                    texts.append("\n".join(chunks))
                if action == "read":
                    return ToolResult(True, "\n---\n".join(texts), data={"path": str(src), "backend": "com"})
                if action == "info":
                    return ToolResult(
                        True,
                        f"path={src}\nbackend=com\nslides={pres.Slides.Count}",
                        data={"path": str(src), "slides": pres.Slides.Count, "backend": "com"},
                    )
                dest = self._resolve(kwargs.get("destination") or str(src))
                dest.parent.mkdir(parents=True, exist_ok=True)
                if action == "append" and kwargs.get("content"):
                    for block in _slides_from_content(kwargs["content"]):
                        slide = pres.Slides.Add(pres.Slides.Count + 1, 12)
                        try:
                            slide.Shapes.Title.TextFrame.TextRange.Text = block.splitlines()[0]
                        except Exception:
                            pass
                elif action in {"write", "save_as"} and kwargs.get("content") is not None:
                    while pres.Slides.Count:
                        pres.Slides(1).Delete()
                    for block in _slides_from_content(kwargs["content"]):
                        slide = pres.Slides.Add(pres.Slides.Count + 1, 12)
                        try:
                            slide.Shapes.Title.TextFrame.TextRange.Text = block.splitlines()[0]
                        except Exception:
                            pass
                pres.SaveAs(str(dest))
                return ToolResult(True, f"Saved {dest}", data={"path": str(dest), "backend": "com"})
            finally:
                pres.Close()

    def _write_com_sheet(self, sheet: Any, table: list[list[str]], start_row: int = 1) -> None:
        for r_idx, row in enumerate(table, start=start_row):
            for c_idx, value in enumerate(row, start=1):
                sheet.Cells(r_idx, c_idx).Value = value
